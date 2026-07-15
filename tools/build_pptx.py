"""Generate the editable PowerPoint version of the Postmeet AI-PM pitch deck.

Run:  python tools/build_pptx.py     (needs: pip install python-pptx)
Out:  docs/postmeet-pitch.pptx

Slide order is deliberate: the PRODUCT is shown first (screenshot, then how it
works) so a reviewer sees what it *is* before any argument about why it matters.
"""
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = pathlib.Path(__file__).resolve().parent.parent
SHOT = ROOT / "docs" / "postmeet-screenshot.png"
OUT = ROOT / "docs" / "postmeet-pitch.pptx"
SHOT_W, SHOT_H = 2368, 1572          # screenshot dimensions (see tools/capture_screenshot.py)

INK   = RGBColor(0x16, 0x18, 0x1F)
MUTED = RGBColor(0x5A, 0x61, 0x6E)
FAINT = RGBColor(0x8B, 0x91, 0x9D)
ACCENT= RGBColor(0xC2, 0x41, 0x0C)
LINE  = RGBColor(0xDE, 0xE1, 0xE7)
SOFT  = RGBColor(0xFB, 0xEF, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG    = RGBColor(0xFB, 0xFB, 0xFC)
BODY, MONO = "Calibri", "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
MW = 13.333


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background(); bg.shadow.inherit = False
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(MW), Inches(0.09))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; rule.line.fill.background(); rule.shadow.inherit = False
    return s


def box(s, l, t, w, h, fill=WHITE, line=LINE, lw=1.0, radius=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if radius:
        try: shp.adjustments[0] = 0.045
        except Exception: pass
    return shp


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.06):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space; p.space_after = Pt(0); p.space_before = Pt(0)
        for (t_, sz, col, bold, fn) in para:
            r = p.add_run(); r.text = t_; f = r.font
            f.size = Pt(sz); f.color.rgb = col; f.bold = bold; f.name = fn
    return tb


def eyebrow(s, text):
    txt(s, 0.9, 0.62, 11.5, 0.4, [[(text, 12.5, ACCENT, True, MONO)]])


def title(s, text, size=33, y=1.15):
    txt(s, 0.9, y, 11.5, 1.5, [[(text, size, INK, True, BODY)]], space=1.02)


# ============ 1 · TITLE ============
s = slide()
txt(s, 0.9, 0.62, 11.5, 0.4, [[("AI PRODUCT CASE STUDY", 12.5, ACCENT, True, MONO)]])
txt(s, 0.9, 1.9, 11.5, 2.2, [
    [("Postmeet", 54, INK, True, BODY), (".", 54, ACCENT, True, BODY)],
    [("The last mile of meetings.", 40, INK, True, BODY)]], space=1.0)
txt(s, 0.9, 4.15, 9.6, 1.2, [[("Every meeting produces commitments. Most of them evaporate. Postmeet turns a raw "
    "transcript into distributed, owned action — in about five seconds, with zero setup.", 17, MUTED, False, BODY)]], space=1.14)
txt(s, 0.9, 5.52, 11.5, 0.4, [[("Shubham Mittal", 15, INK, True, BODY), ("   ·   AI Product Manager", 15, MUTED, False, BODY)]])
# Award badge — a real, verifiable credential; earns attention before any claim does.
box(s, 0.9, 6.0, 8.9, 0.44, fill=SOFT, line=ACCENT, lw=1.2, radius=True)
txt(s, 1.12, 6.06, 8.6, 0.34, [[("🏆  Top 10 Finalist → 2nd Runner-Up", 11.5, ACCENT, True, MONO),
    ("  ·  PromptWars Chennai 2026  ·  Google for Developers × Hack2Skill", 10, MUTED, False, MONO)]])
txt(s, 0.9, 6.62, 11.5, 0.4, [[("Live demo: postmeet.onrender.com      Repo: github.com/Shubham-33/postmeet", 12.5, FAINT, False, MONO)]])

# ============ 2 · SEE IT WORK (product first — the visual hook) ============
s = slide(); eyebrow(s, "—  SEE IT WORK")
title(s, "Paste a transcript → the board writes itself.", size=28)
_ih = 5.15; _iw = _ih * SHOT_W / SHOT_H; _left = (MW - _iw) / 2
s.shapes.add_picture(str(SHOT), Inches(_left), Inches(1.8), height=Inches(_ih))
box(s, _left, 1.8, _iw, _ih, fill=None, line=LINE, lw=1.0)
txt(s, 0.9, 7.02, 11.5, 0.35, [[("Real output from the live app — extracted in ~3s: summary, decisions, and one "
    "column per owner, each card one click from a prefilled Calendar event or Gmail draft.", 10.5, FAINT, False, BODY)]])

# ============ 3 · THE PRODUCT ============
s = slide(); eyebrow(s, "01 · THE PRODUCT")
title(s, "Paste a meeting → get owned action → distribute in one click.")
steps = [("INPUT", "Drop it in", "Paste a transcript, a public Google Doc URL, or a file. No account."),
         ("EXTRACT", "AI structures it", "Summary, decisions, and action items — each with owner, email, due date, context."),
         ("DISTRIBUTE", "One-click send", "Each item opens a prefilled Google Calendar event or Gmail draft. Just Save / Send.")]
x = 0.9; w = 3.75
for k, h, b in steps:
    box(s, x, 2.55, w, 2.0, fill=WHITE, line=LINE, radius=True)
    txt(s, x + 0.28, 2.75, w - 0.5, 0.4, [[(k, 11, ACCENT, True, MONO)]])
    txt(s, x + 0.28, 3.15, w - 0.5, 0.5, [[(h, 17, INK, True, BODY)]])
    txt(s, x + 0.28, 3.7, w - 0.5, 0.8, [[(b, 13, MUTED, False, BODY)]], space=1.12)
    x += w + 0.2
box(s, 0.9, 5.05, 0.06, 1.35, fill=ACCENT, line=None)
txt(s, 1.2, 5.12, 11.2, 1.3, [[("The “no-OAuth trick”: prefilled Calendar / Gmail URLs open in the user's "
    "already-logged-in Google. ", 18, INK, False, BODY),
    ("Time-to-value ~10 seconds — no integration, no admin approval, no scope review.", 18, ACCENT, True, BODY)]], space=1.14)

# ============ 4 · THE PROBLEM ============
s = slide(); eyebrow(s, "02 · THE PROBLEM")
title(s, "Meetings manufacture commitments. Then the commitments go missing.")
txt(s, 0.9, 2.5, 4.4, 1.4, [[("~70%", 66, ACCENT, True, BODY)]])
txt(s, 0.9, 3.9, 4.9, 1.9, [[("of meeting action items are never completed — they live in someone's notes, "
    "a Slack thread, a doc nobody reopens.", 15.5, MUTED, False, BODY)]], space=1.16)
pts = [("who", "Knowledge workers spend a huge share of the week in meetings — what matters is the follow-through, not the recording."),
       ("gap", "The pain isn't capturing what was said. It's doing what was decided."),
       ("cost", "Dropped commitments = re-work, missed deadlines, and eroded trust.")]
y = 2.55
for k, v in pts:
    txt(s, 6.0, y, 1.0, 0.4, [[(k, 12.5, ACCENT, True, MONO)]])
    txt(s, 7.0, y - 0.02, 5.4, 1.1, [[(v, 14.5, MUTED, False, BODY)]], space=1.12); y += 1.15
txt(s, 0.9, 6.7, 11.5, 0.5, [[("Industry estimates range ~44–73% never completed (Fellow, Streamli9). "
    "Directional — as PM, first job is to validate with first-party instrumentation before over-investing.", 11.5, FAINT, False, BODY)]], space=1.1)

# ============ 5 · THE INSIGHT ============
s = slide(); eyebrow(s, "03 · THE INSIGHT")
title(s, "Everyone is fighting over capture. The last mile is wide open.")
box(s, 0.9, 2.5, 5.75, 2.35, fill=WHITE, line=LINE, radius=True)
txt(s, 1.2, 2.72, 5.2, 0.4, [[("CROWDED — “CAPTURE”", 11, FAINT, True, MONO)]])
txt(s, 1.2, 3.18, 5.2, 0.5, [[("Transcribe & summarize", 19, INK, True, BODY)]])
txt(s, 1.2, 3.75, 5.2, 1.0, [[("Otter, Fireflies, Fathom, Granola — all competing on better transcripts. Table "
    "stakes now. The commitments they surface still die in a doc.", 13.5, MUTED, False, BODY)]], space=1.12)
box(s, 6.85, 2.5, 5.6, 2.35, fill=SOFT, line=ACCENT, lw=1.4, radius=True)
txt(s, 7.15, 2.72, 5.0, 0.4, [[("OPEN — “THE LAST MILE”", 11, ACCENT, True, MONO)]])
txt(s, 7.15, 3.18, 5.0, 0.5, [[("Distribute & act", 19, INK, True, BODY)]])
txt(s, 7.15, 3.75, 5.0, 1.0, [[("Getting each commitment into the tool the owner actually lives in — their calendar, "
    "their inbox — is where follow-through is won or lost.", 13.5, MUTED, False, BODY)]], space=1.12)
box(s, 0.9, 5.35, 0.06, 1.2, fill=ACCENT, line=None)
txt(s, 1.2, 5.45, 11.2, 1.1, [[("The blocker on the last mile isn't a smarter model — it's ", 21, INK, False, BODY),
    ("setup friction", 21, ACCENT, True, BODY), (". So the wedge is zero-setup distribution.", 21, INK, False, BODY)]], space=1.14)

# ============ 6 · THE AI UNDER THE HOOD ============
s = slide(); eyebrow(s, "04 · THE AI UNDER THE HOOD")
title(s, "Reliable structure out of an unreliable medium.")
txt(s, 0.9, 2.35, 11.5, 0.5, [[("transcript  →  ", 13, MUTED, False, MONO), ("LLM · JSON mode", 13, ACCENT, True, MONO),
    ("  →  schema-validated JSON  →  board + prefilled links", 13, MUTED, False, MONO)]])
cards = [("Structured output", "A fixed JSON contract (summary / decisions / action items), pinned in the system prompt and decoded by a tolerant parser. One call, no retry loop."),
         ("Guardrails", "Extracts only explicit commitments — no invented tasks. Empty fields when data is absent. The difference between useful and hallucinated."),
         ("Graceful degradation", "A model fallback chain: if the primary is cold or errors, it drops to a backup so the request still succeeds.")]
x = 0.9; w = 3.75
for h, b in cards:
    box(s, x, 3.15, w, 2.6, fill=WHITE, line=LINE, radius=True)
    txt(s, x + 0.28, 3.4, w - 0.5, 0.5, [[(h, 16.5, INK, True, BODY)]])
    txt(s, x + 0.28, 4.0, w - 0.5, 1.6, [[(b, 13, MUTED, False, BODY)]], space=1.14)
    x += w + 0.2

# ============ 7 · THE DEFINING DECISION ============
s = slide(); eyebrow(s, "05 · THE DEFINING DECISION")
title(s, "I shipped the smaller model on purpose.")


def optbox(x, tag, tagcol, name, rows, win):
    box(s, x, 2.4, 5.6, 2.15, fill=(SOFT if win else WHITE), line=(ACCENT if win else LINE),
        lw=(1.4 if win else 1.0), radius=True)
    txt(s, x + 0.3, 2.6, 5.0, 0.35, [[(tag, 10.5, tagcol, True, MONO)]])
    txt(s, x + 0.3, 2.98, 5.0, 0.5, [[(name, 20, INK, True, MONO)]])
    yy = 3.6
    for k, v in rows:
        txt(s, x + 0.3, yy, 1.6, 0.35, [[(k, 11.5, FAINT, False, MONO)]])
        txt(s, x + 1.9, yy, 3.4, 0.35, [[(v, 13, INK, True, BODY)]], align=PP_ALIGN.RIGHT); yy += 0.34


optbox(0.9, "REJECTED AS PRIMARY", FAINT, "Llama 3.1 70B",
       [("latency", "9s–46s, inconsistent"), ("reliability", "cold-starts, timeouts"), ("quality", "marginally higher")], False)
optbox(6.85, "SHIPPED AS PRIMARY", ACCENT, "Llama 3.1 8B",
       [("latency", "~2.5s, consistent"), ("reliability", "steady across runs"), ("quality", "comparable")], True)
box(s, 0.9, 5.05, 0.06, 1.5, fill=ACCENT, line=None)
txt(s, 1.2, 5.12, 11.2, 1.5, [[("Model quality is an ", 18, INK, False, BODY), ("input", 18, ACCENT, True, BODY),
    (" to product value, not the goal. A demo that always works in two seconds beats a “smarter” one "
     "that sometimes hangs — reliability builds trust faster than marginal accuracy.", 18, INK, False, BODY)]], space=1.16)

# ============ 8 · METRICS ============
s = slide(); eyebrow(s, "06 · HOW I'D MEASURE SUCCESS")
title(s, "One north star, a funnel beneath it, guardrails around it.")
box(s, 0.9, 2.4, 11.55, 1.25, fill=SOFT, line=ACCENT, lw=1.4, radius=True)
txt(s, 1.2, 2.55, 11.0, 0.35, [[("NORTH-STAR METRIC", 11, ACCENT, True, MONO)]])
txt(s, 1.2, 2.9, 11.0, 0.5, [[("Commitment Completion Rate", 21, INK, True, BODY)]])
txt(s, 1.2, 3.38, 11.0, 0.3, [[("Of the action items surfaced, what share actually get done? Everything else is a proxy.", 13, MUTED, False, BODY)]])
txt(s, 0.9, 3.95, 6.4, 0.4, [[("INPUT METRICS — THE FUNNEL", 11, FAINT, True, MONO)]])
for i, c in enumerate(["Activation: % sessions that extract → distribute ≥1 item",
                       "Time-to-value: seconds paste → first send",
                       "Extraction quality: precision / recall vs. eval set",
                       "Retention: repeat use / week"]):
    txt(s, 0.9, 4.35 + i * 0.42, 6.2, 0.4, [[("•  ", 13, ACCENT, True, BODY), (c, 13.5, MUTED, False, BODY)]])
txt(s, 7.6, 3.95, 4.8, 0.4, [[("GUARDRAIL METRICS", 11, FAINT, True, MONO)]])
for i, c in enumerate(["Hallucination rate (invented items)", "Edit rate (proxy for accuracy)", "Error / failure rate"]):
    txt(s, 7.6, 4.35 + i * 0.42, 4.8, 0.4, [[("•  ", 13, ACCENT, True, BODY), (c, 13.5, MUTED, False, BODY)]])

# ============ 9 · EVALS ============
s = slide(); eyebrow(s, "07 · MAKING AI DECISIONS RIGOROUSLY")
title(s, "Vibes don't scale. Evals do.")
pts = [("set", "Build a golden eval set: labelled transcripts with the “true” decisions and action items."),
       ("score", "Measure precision (did we invent tasks?), recall (did we miss commitments?), and field accuracy."),
       ("use", "That's how I'd choose 8B vs 70B, catch regressions, and set a launch quality bar — not by gut.")]
y = 2.6
for k, v in pts:
    txt(s, 0.9, y, 1.0, 0.4, [[(k, 12.5, ACCENT, True, MONO)]])
    txt(s, 1.9, y - 0.02, 4.9, 1.0, [[(v, 15, MUTED, False, BODY)]], space=1.14); y += 1.2
box(s, 7.15, 2.55, 5.3, 3.5, fill=WHITE, line=LINE, radius=True)
txt(s, 7.45, 2.78, 4.7, 0.35, [[("THE COMPOUNDING LOOP", 11, ACCENT, True, MONO)]])
txt(s, 7.45, 3.18, 4.7, 0.5, [[("Every edit is a label", 18, INK, True, BODY)]])
txt(s, 7.45, 3.8, 4.7, 2.0, [[("Users correct the extracted items before sending. Those corrections are free, "
    "in-domain training data → fine-tune on accepted output → org-specific accuracy competitors can't "
    "match without the same data. The eval loop becomes the moat.", 14, MUTED, False, BODY)]], space=1.16)

# ============ 10 · PRIORITIZATION ============
s = slide(); eyebrow(s, "08 · PRIORITIZATION")
title(s, "What I deliberately left out — and why that's the point.")
txt(s, 0.9, 2.2, 11.4, 0.7, [[("MVP scope is a hypothesis test, not a feature checklist. I cut everything not needed to validate "
    "one thing: will people extract, then actually distribute?", 15.5, MUTED, False, BODY)]], space=1.14)
items = [("CUT · PERSISTENCE", "No database. A refresh clears everything — and “no data stored” removes privacy friction for testers."),
         ("CUT · REAL INTEGRATIONS", "The URL-prefill trick proves value without OAuth. Real APIs are a scaling investment, made once the loop is proven."),
         ("CUT · ACCOUNTS & TEAMS", "Single-session, no login. Multi-user is a growth concern, not a validation one."),
         ("KEPT · THE CORE LOOP + TRUST", "Extraction reliability, one-click distribution, honest guardrails, 100%-tested logic.")]
xs = [0.9, 6.85]; ys = [3.2, 4.95]
for i, (h, b) in enumerate(items):
    x = xs[i % 2]; y = ys[i // 2]
    box(s, x, y, 5.6, 1.6, fill=WHITE, line=LINE, radius=True)
    txt(s, x + 0.28, y + 0.18, 5.0, 0.35, [[(h, 11, ACCENT, True, MONO)]])
    txt(s, x + 0.28, y + 0.55, 5.05, 0.95, [[(b, 13, MUTED, False, BODY)]], space=1.12)

# ============ 11 · ROADMAP ============
s = slide(); eyebrow(s, "09 · ROADMAP & MOAT")
title(s, "From a one-click helper to the system of record for follow-through.")
hor = [("NOW · VALIDATE", "The core loop", ["Paste → extract → distribute", "Zero setup, no OAuth", "Reliable structured output"], True),
       ("NEXT · AUTO-CAPTURE", "Capture everywhere", ["Extension + AI notetaker", "Auto-joins Zoom / Teams / Meet", "Native Gmail, Outlook, Calendar"], False),
       ("LATER · OWN THE OUTCOME", "Proactive follow-through", ["Persistence + team workspaces", "Nudges: “did you ship X?”", "Completion tracking → closes the loop"], False)]
x = 0.9; w = 3.75
for k, h, rows, now in hor:
    box(s, x, 2.5, w, 2.6, fill=WHITE, line=(ACCENT if now else LINE), lw=(1.4 if now else 1.0), radius=True)
    txt(s, x + 0.26, 2.7, w - 0.45, 0.35, [[(k, 10.5, ACCENT, True, MONO)]])
    txt(s, x + 0.26, 3.08, w - 0.45, 0.45, [[(h, 15.5, INK, True, BODY)]])
    yy = 3.62
    for r in rows:
        txt(s, x + 0.26, yy, w - 0.45, 0.5, [[("•  ", 12, ACCENT, True, BODY), (r, 12.5, MUTED, False, BODY)]], space=1.05); yy += 0.46
    x += w + 0.2
box(s, 0.9, 5.4, 0.06, 1.1, fill=ACCENT, line=None)
txt(s, 1.2, 5.46, 11.2, 1.1, [[("Moat = the feedback loop. Accepted & edited action items become proprietary training "
    "data → accuracy that compounds with usage.", 17, INK, False, BODY)]], space=1.14)

# ============ 12 · CLOSE ============
s = slide(); eyebrow(s, "10 · WHY THIS MATTERS FOR THE ROLE")
title(s, "I don't just use AI tools. I make product decisions about them.", size=30)
comp = [("Problem framing", "Found the non-obvious wedge — the last mile, not capture; friction, not model size."),
        ("AI judgment", "Model tradeoffs, evals, guardrails, graceful degradation — decisions, with reasons."),
        ("Metrics rigor", "A north star, an input funnel, and guardrails — measurement designed in."),
        ("Prioritization", "Scoped an MVP as a hypothesis test; cut with intent; shipped what de-risks it."),
        ("Vision", "A credible path from a one-click tool to a defensible product with a data moat."),
        ("End-to-end", "Built, tested (100% coverage), and deployed it solo — insight to live URL.")]
xs = [0.9, 5.08, 9.26]; w = 3.95
for i, (h, b) in enumerate(comp):
    x = xs[i % 3]; y = 2.65 if i < 3 else 4.15
    txt(s, x, y, w, 0.4, [[(h, 15, INK, True, BODY)]])
    txt(s, x, y + 0.42, w, 1.1, [[(b, 12.5, MUTED, False, BODY)]], space=1.12)
txt(s, 0.9, 5.9, 11.5, 0.4, [[("Shubham Mittal", 14, INK, True, BODY),
    ("   ·   postmeet.onrender.com   ·   shubham-33.github.io/postmeet/case-study/", 12.5, FAINT, False, MONO)]])

prs.save(str(OUT))
print(f"saved {OUT}  ·  slides: {len(prs.slides._sldIdLst)}")
